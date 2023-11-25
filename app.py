import subprocess
import brotli
import pdfplumber
import pandas as pd
import tabula
from pdfminer.high_level import extract_text
from flask import Flask, render_template, request, redirect, url_for, send_from_directory, jsonify, send_file
import os
import io
import traceback
import shutil
import zipfile
from PIL import Image
from reportlab.lib.pagesizes import letter, landscape, A4
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, PageBreak
from openpyxl import load_workbook
from reportlab.pdfgen import canvas
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, PageBreak
import csv
from PyPDF2 import PdfReader, PdfWriter, PdfFileMerger, PageObject
from werkzeug.utils import secure_filename
from pdf2docx import Converter
import PyPDF2
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_path
from PyPDF2 import PdfReader
from PyPDF2 import PdfMerger
import PyPDF2
# Naming current module
app = Flask(__name__)

# Define upload and converted PDF directories
UPLOAD_FOLDER = 'uploads'
CONVERTED_FOLDER = 'converted_pdfs'
TEMP_FOLDER = 'tmp'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['CONVERTED_FOLDER'] = CONVERTED_FOLDER
app.config['TEMP_FOLDER'] = TEMP_FOLDER

# Render the home page for file conversion
@app.route('/')
def index():
    return render_template('index.html')

# Function to convert images to PDF
def convert_image_to_pdf(image_path, output_path):
    # A4 size
    pdf_canvas = canvas.Canvas(output_path, pagesize=(595, 842))
    img = Image.open(image_path)
    img_width, img_height = img.size
    scale_factor = min(595 / img_width, 842 / img_height)
    x_offset = (595 - img_width * scale_factor) / 2
    y_offset = (842 - img_height * scale_factor) / 2
    # Draw the image on the PDF
    pdf_canvas.drawInlineImage(image_path, x_offset, y_offset, width=img_width * scale_factor, height=img_height * scale_factor)
    # Save the PDF
    pdf_canvas.save()


# Specify the path to the LibreOffice binary
LIBREOFFICE_PATH = '/usr/bin'

# Set the UNO_PATH environment variable
os.environ['UNO_PATH'] = LIBREOFFICE_PATH

# Function to convert document types to PDF
def convert_doc_to_pdf(document_path, output_path):
    try:
        # Using libreoffice to convert the document to PDF
        subprocess_args = [
            'libreoffice',
            '--headless',
            '--convert-to', 'pdf:writer_pdf_Export',
            '--outdir', os.path.dirname(output_path),
            '--writer_pdf_Export_PageSize', 'A4',
            document_path
        ]
        subprocess.run(subprocess_args, check=True)
        return True
    except subprocess.CalledProcessError as e:
        print(f"Conversion failed: {e}")
        return False



def convert_ppt_to_pdf(ppt_path, output_path):
    try:
        # Using libre office for converting ppt to PDF
        env = os.environ.copy()
        env['PDFA1B_OUTDIR'] = os.path.dirname(output_path)
        subprocess_args = [os.path.join(LIBREOFFICE_PATH, 'libreoffice'), '--headless', '--convert-to', 'pdf:writer_pdf_Export', ppt_path]
        subprocess.run(subprocess_args, check=True, env=env)
        return True
    except subprocess.CalledProcessError as e:
        print(f"Conversion failed: {e}")
        return False


# Function to convert CSV to PDF
def convert_csv_to_pdf(csv_path, pdf_path):
    # Reading the CSV file
    data = []
    with open(csv_path, 'r') as csvfile:
        csvreader = csv.reader(csvfile)
        for row in csvreader:
            data.append(row)

    # Create a PDF template using reportlab
    doc = SimpleDocTemplate(pdf_path, pagesize=A4)
    elements = []

    font_name = 'Helvetica-Bold'
    font_size = 8

    columns_per_table = 5

    num_columns = len(data[0])

    for start_col in range(0, num_columns, columns_per_table):
        end_col = start_col + columns_per_table
        table_data = [row[start_col:end_col] for row in data]

        # Generate table with data for PDF
        table = Table(table_data)

        col_widths = [1.5] * len(table_data[0])
        style = TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), (0.9, 0.9, 0.9)),
            ('TEXTCOLOR', (0, 0), (-1, 0), (0, 0, 0)),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), font_name),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), (0.85, 0.85, 0.85)),
            ('GRID', (0, 0), (-1, -1), 1, (0, 0, 0)),
            ('FONTSIZE', (0, 0), (-1, -1), font_size),
        ])
        table.setStyle(style)
        elements.append(table)

        if end_col < num_columns:
            elements.append(PageBreak())

    doc.build(elements)

# Function to convert Excel to PDF
def convert_excel_to_pdf(excel_path, output_path):
    # Create a landscape PDF using the reportlab library
    doc = SimpleDocTemplate(output_path, pagesize=A4)
    elements = []
    workbook = load_workbook(excel_path)

    # Font size and style
    font_name = 'Helvetica-Bold'
    font_size = 12

    # Limiting colunm in pdf
    columns_per_table = 5

    # Reading Sheets from Excel
    for sheet in workbook.sheetnames:
        data = []
        for row in workbook[sheet].iter_rows(values_only=True):
            data.append(row)

        num_columns = len(data[0])

        for start_col in range(0, num_columns, columns_per_table):
            end_col = start_col + columns_per_table
            table_data = [row[start_col:end_col] for row in data]

            # Generate table with data for PDF
            table = Table(table_data)

            col_widths = [1.5] * len(table_data[0])
            style = TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), (0.8, 0.8, 0.8)),
                ('TEXTCOLOR', (0, 0), (-1, 0), (1, 1, 1)),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), font_name),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), (0.95, 0.95, 0.95)),
                ('FONTSIZE', (0, 0), (-1, -1), font_size),
            ])
            table.setStyle(style)
            elements.append(table)
            if end_col < num_columns:
                elements.append(PageBreak())
    doc.build(elements)


# Function converting PDF to image
def convert_pdf_to_image(pdf_path, output_dir):
    images = convert_from_path(pdf_path)
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    image_filenames = []
    for i, image in enumerate(images):
        image_filename = f"page_{i + 1}.png"
        image.save(os.path.join(output_dir, image_filename))
        image_filenames.append(image_filename)
    return image_filenames

#Return the images from directory
@app.route('/converted_images/<filename>')
def converted_images(filename):
    return send_from_directory(app.config['CONVERTED_FOLDER'], filename)

#Converting PDF to Excel file
def convert_pdf_to_xlsx(pdf_path, xlsx_path):
    # Open the PDF file
    with pdfplumber.open(pdf_path) as pdf:
        all_tables = []
        for page in pdf.pages:
            # Extract tables from the page
            tables = page.extract_tables()
            all_tables.extend(tables)
        # Using pandans dataframe for reading table
        df = pd.DataFrame([row for table in all_tables for row in table])
        # Save the dataframe to an excel file
        df.to_excel(xlsx_path, index=False)

#Converting PDF to Csv
def convert_pdf_to_csv(pdf_path, csv_path):
    tabula.convert_into(pdf_path, csv_path, output_format="csv")

@app.route('/converted_files/<filename>')
def converted_files(filename):
    return send_from_directory(app.config['CONVERTED_FOLDER'], filename)

#Converting PDF to text
def convert_pdf_to_text(pdf_path, output_path):
    text = extract_text(pdf_path)
    with open(output_path, "w", encoding="utf-8") as text_file:
        text_file.write(text)


# Function to convert PDF to DOCX
def convert_pdf_to_docx(pdf_path, docx_path):
    try:
        cv = Converter(pdf_path)
        cv.convert(docx_path, start=0, end=None)
        cv.close()
        return True
    except Exception as e:
        print(f"Conversion failed: {e}")
        return False

# Function to convert PDF to PPTX
def convert_pdf_to_pptx(pdf_path, pptx_path):
    # Create a new PowerPoint presentation
    prs = Presentation()
    with open(pdf_path, 'rb') as pdf_file:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        for page_num, page in enumerate(pdf_reader.pages):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            shapes = slide.shapes
            title_shape = shapes.title
            if len(shapes.placeholders) > 1:
                body_shape = shapes.placeholders[1]
            else:
                body_shape = shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
            title_shape.text = f"Slide {page_num + 1}"
            tf = body_shape.text_frame
            text = page.extract_text()
            text = text.replace('\n', ' ')
            paragraphs = text.split('\r')
            for paragraph in paragraphs:
                p = tf.add_paragraph()
                p.text = paragraph
    prs.save(pptx_path)


# Rendering the merging function
@app.route('/merge', methods=['GET', 'POST'])
def merge_pdfs():
    if request.method == 'POST':
        pdf_files = request.files.getlist('pdf_files')
        if pdf_files:
            merged_pdf_path = os.path.join(app.config['CONVERTED_FOLDER'], 'merged.pdf')
            merge_pdfs([pdf.filename for pdf in pdf_files], merged_pdf_path)
            merged_pdf_url = url_for('uploaded_pdf', filename='merged.pdf')
            return render_template('merge.html', merged_pdf_url=merged_pdf_url)
        else:
            return "No PDF files to merge."

    return render_template('merge.html')

# Function to merge PDF files
def merge_pdfs(pdf_files, output_path):
    pdf_merger = PdfMerger()

    for pdf_file in pdf_files:
        pdf_merger.append(os.path.join(app.config['UPLOAD_FOLDER'], pdf_file))

    pdf_merger.write(output_path)
    pdf_merger.close()


os.makedirs(app.config['CONVERTED_FOLDER'], exist_ok=True)

# Function to create a zip file
def create_zip_file(zip_filename, files):
    with zipfile.ZipFile(zip_filename, 'w') as zip_file:
        for file in files:
            zip_file.write(file, os.path.basename(file))

# Route to download all split PDFs as a zip file
@app.route('/download_all_zip', methods=['GET', 'POST'])
def download_all_zip():
    if request.method == 'POST':
        pdf_file = request.files['pdf_file']
        if pdf_file:
            output_dir = app.config['CONVERTED_FOLDER']
            os.makedirs(output_dir, exist_ok=True)
            # Define the pdf_reader within the function
            pdf_data = pdf_file.read()
            pdf_reader = PdfReader(io.BytesIO(pdf_data))
            split_pdf_paths = []
            for page_num, page in enumerate(pdf_reader.pages):
                pdf_writer = PdfWriter()
                pdf_writer.add_page(page)
                split_page_path = os.path.join(output_dir, f'page_{page_num + 1}.pdf')
                split_pdf_paths.append(split_page_path)

                with open(split_page_path, 'wb') as output_pdf:
                    pdf_writer.write(output_pdf)

            # Create a zip file
            zip_filename = os.path.join(output_dir, 'split_pdf.zip')
            create_zip_file(zip_filename, split_pdf_paths)

            # Return the zip file and split PDF paths for download
            return render_template('split.html', split_pdf_paths=split_pdf_paths, zip_filename=zip_filename)

    return render_template('split.html')

# Route to download a split PDF
@app.route('/download/<path:filename>')
def download_split_pdf(filename):
    return send_file(filename, as_attachment=True)
# Route to download the zip file
@app.route('/download_zip/<path:zip_filename>')
def download_zip(zip_filename):
    return send_file(zip_filename, as_attachment=True)



# Rendering the locking function
@app.route('/lock', methods=['GET', 'POST'])
def lock_pdf():
    if request.method == 'POST':
        pdf_file = request.files['pdf_file']
        password = request.form['password']
        confirm_password = request.form['confirm_password']
        if pdf_file and password == confirm_password:
            output_dir = app.config['CONVERTED_FOLDER']
            os.makedirs(output_dir, exist_ok=True)


            locked_pdf_path = os.path.join(output_dir, secure_filename(pdf_file.filename))
            pdf_file.save(locked_pdf_path)

            set_pdf_password(locked_pdf_path, password)

            locked_pdf_url = url_for('uploaded_pdf', filename=os.path.basename(locked_pdf_path))
            return render_template('lock.html', locked_pdf_url=locked_pdf_url)
        else:
            return "Passwords do not match. Please try again."

    return render_template('lock.html')

# Function to set a password for a PDF
def set_pdf_password(pdf_path, password):
    pdf_writer = PdfWriter()

    with open(pdf_path, 'rb') as pdf_file:
        pdf_reader = PdfReader(pdf_file)

        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            pdf_writer.add_page(page)

        pdf_writer.encrypt(password)

    with open(pdf_path, 'wb') as output_pdf:
        pdf_writer.write(output_pdf)


# Rendering the unlocking function
@app.route('/unlock', methods=['GET', 'POST'])
def unlock_pdf():
    if request.method == 'POST':
        pdf_file = request.files['pdf_file']
        password = request.form['password']

        temp_folder = app.config['TEMP_FOLDER']
        temp_pdf_path = os.path.join(temp_folder, secure_filename(pdf_file.filename))
        pdf_file.save(temp_pdf_path)

        if is_pdf_locked(temp_pdf_path, password):

            with open(temp_pdf_path, 'rb') as temp_file:
                pdf_reader = PdfReader(temp_file)
                pdf_reader.decrypt(password)
                pdf_writer = PdfWriter()
                for page_num in range(len(pdf_reader.pages)):
                    pdf_writer.add_page(pdf_reader.pages[page_num])

                unlocked_pdf_path = os.path.join(app.config['CONVERTED_FOLDER'], secure_filename(pdf_file.filename))
                with open(unlocked_pdf_path, 'wb') as unlocked_file:
                    pdf_writer.write(unlocked_file)
            os.remove(temp_pdf_path)

            unlocked_pdf_url = url_for('uploaded_pdf', filename=os.path.basename(unlocked_pdf_path))
            return render_template('unlock.html', unlocked_pdf_url=unlocked_pdf_url)
        else:

            os.remove(temp_pdf_path)
            return "Incorrect password. Please try again."

    return render_template('unlock.html')

os.makedirs('converted_pdfs', exist_ok=True)


def compress_file(input_path, output_path, compression_level):
    try:
        with open(input_path, 'rb') as input_file:
            pdf_reader = PdfReader(input_file)
            pdf_writer = PdfWriter()

            for page_num in range(pdf_reader.getNumPages()):
                page = pdf_reader.getPage(page_num)
                content = page.compressContentStreams()

                if content is not None:
                    # Compress the content using brotli library
                    compressed_content = brotli.compress(content, quality=compression_level)
                    page.__setitem__(PageObject.CONTENTS, compressed_content)

                    pdf_writer.addPage(page)

            with open(output_path, 'wb') as output_file:
                pdf_writer.write(output_file)

        print(f"Compression successful: {output_path}")
    except Exception as e:
        print(f"Error during compression: {e}")


@app.route('/compress', methods=['GET', 'POST'])
def compress():
    import os
    if request.method == 'POST':
        compression_level = int(request.form['compression_level'])
        file = request.files['file']

        if file:
            # Save the uploaded file
            input_path = os.path.join('output', file.filename)
            file.save(input_path)

            # Compress the file
            output_path = os.path.join('output', f'compressed_{file.filename}.zip')
            compress_file(input_path, output_path, compression_level)

            return render_template('compress.html', compressed_file=output_path, os=os)

    return render_template('compress.html', compressed_file=None, os=os)

@app.route('/download_compressed_pdf/<filename>')
def download_compressed_pdf(filename):
    compressed_file_path = os.path.join(app.config['CONVERTED_FOLDER'], filename)
    return send_file(compressed_file_path, as_attachment=True)

# Function to check if a PDF is locked
def is_pdf_locked(pdf_path, password):
    with open(pdf_path, 'rb') as pdf_file:
        pdf_reader = PdfReader(pdf_file)
        return pdf_reader.is_encrypted and pdf_reader.decrypt(password)

# Rendering the delete pages function
@app.route('/delete', methods=['GET', 'POST'])
def delete_pages():
    if request.method == 'POST':
        pdf_file = request.files['pdf_file']
        pages_to_delete = request.form['pages_to_delete']
        # Save the uploaded PDF to a temporary location
        temp_pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], pdf_file.filename)
        pdf_file.save(temp_pdf_path)

        output_pdf_path = os.path.join(app.config['CONVERTED_FOLDER'], 'pages_deleted.pdf')
        delete_pdf_pages(temp_pdf_path, output_pdf_path, pages_to_delete)

        deleted_pdf_url = url_for('uploaded_pdf', filename='pages_deleted.pdf')
        return render_template('delete.html', deleted_pdf_url=deleted_pdf_url)

    return render_template('delete.html')

# Function to delete specific pages from a PDF
def delete_pdf_pages(input_path, output_path, pages_to_delete):
    pdf_reader = PdfReader(input_path)
    pdf_writer = PdfWriter()
    for page_num in range(len(pdf_reader.pages)):
        if str(page_num + 1) not in pages_to_delete.split(','):
            pdf_writer.add_page(pdf_reader.pages[page_num])

    with open(output_path, 'wb') as output_pdf:
        pdf_writer.write(output_pdf)


# Rendering the repair function
@app.route('/repair', methods=['GET', 'POST'])
def repair_pdf():
    if request.method == 'POST':
        pdf_file = request.files['pdf_file']
        temp_pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], pdf_file.filename)
        pdf_file.save(temp_pdf_path)

        output_pdf_path = os.path.join(app.config['CONVERTED_FOLDER'], 'repaired.pdf')
        repair_pdf_file(temp_pdf_path, output_pdf_path)

        repaired_pdf_url = url_for('uploaded_pdf', filename='repaired.pdf')
        return render_template('repair.html', repaired_pdf_url=repaired_pdf_url)

    return render_template('repair.html')

# Function to repair a PDF using PyPDF2
def repair_pdf_file(input_path, output_path):
    pdf_reader = PdfReader(input_path)
    pdf_writer = PdfWriter()

    for page_num in range(pdf_reader.getNumPages()):
        page = pdf_reader.getPage(page_num)
        pdf_writer.addPage(page)

    with open(output_path, 'wb') as output_pdf:
        pdf_writer.write(output_pdf)

@app.route('/rotate', methods=['GET', 'POST'])
def rotate_pdf():
    if request.method == 'POST':
        pdf_file = request.files['pdf_file']
        rotation_direction = request.form.get('rotate_direction')

        if not rotation_direction:
            return "Rotation direction not specified."

        if pdf_file:
            temp_folder = app.config['TEMP_FOLDER']
            temp_pdf_path = os.path.join(temp_folder, secure_filename(pdf_file.filename))
            pdf_file.save(temp_pdf_path)

            rotated_pdf_path = os.path.join(app.config['CONVERTED_FOLDER'], 'rotated.pdf')
            rotate_pdf_pages(temp_pdf_path, rotated_pdf_path, rotation_direction)

            os.remove(temp_pdf_path)

            rotated_pdf_url = url_for('uploaded_pdf', filename='rotated.pdf')
            return render_template('rotate.html', rotated_pdf_url=rotated_pdf_url)

    return render_template('rotate.html')

def rotate_pdf_pages(input_path, output_path, rotation_direction):
    with open(input_path, 'rb') as input_file, open(output_path, 'wb') as output_file:
        pdf_reader = PdfReader(input_file)
        pdf_writer = PdfWriter()

        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            rotation_angle = 90 if rotation_direction == 'right' else -90
            page.rotate(rotation_angle)
            pdf_writer.add_page(page)
        pdf_writer.write(output_file)


# Rendering the converting function for converting supported formats
@app.route('/convert', methods=['POST'])
def convert_file():
    if 'file' not in request.files:
        return redirect(request.url)
    file = request.files['file']
    if file.filename == '':
        return redirect(request.url)
    if file:
        filename = file.filename
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        pdf_filename = f'{os.path.splitext(filename)[0]}.pdf'
        supported_extensions = {
            '.pptx', '.ppt', '.pps', '.ppsx', '.odp',
            '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.pbm', '.pgm', '.ppm', '.xbm', '.webp',
            '.doc', '.docx', '.docm', '.dot', '.dotm', '.dotx', '.odt', '.rtf', '.txt', '.wps', '.xml', '.xps',
            '.xlsx', '.xlsm', '.xltx', '.csv', '.pdf', '.odt', '.ods', '.odp', '.odg', '.odf'
        }
        file_extension = os.path.splitext(filename)[-1].lower()
        if file_extension in supported_extensions:
            if file_extension == '.pdf':
                conversion_type = request.form['conversion_type']

                if conversion_type == 'pdf_to_image':
                    image_filenames = convert_pdf_to_image(file_path, app.config['CONVERTED_FOLDER'])
                    image_urls = [url_for('converted_images', filename=image_filename) for image_filename in
                                  image_filenames]
                    return render_template('index.html', image_urls=image_urls)
                elif conversion_type == 'pdf_to_xlsx':
                    xlsx_filename = f'{os.path.splitext(filename)[0]}.xlsx'
                    convert_pdf_to_xlsx(file_path, os.path.join(app.config['CONVERTED_FOLDER'], xlsx_filename))
                    xlsx_url = url_for('converted_files', filename=xlsx_filename)
                    return render_template('index.html', xlsx_url=xlsx_url)
                elif conversion_type == 'pdf_to_csv':
                    csv_filename = f'{os.path.splitext(filename)[0]}.csv'
                    convert_pdf_to_csv(file_path, os.path.join(app.config['CONVERTED_FOLDER'], csv_filename))
                    csv_url = url_for('converted_files', filename=csv_filename)
                    return render_template('index.html', csv_url=csv_url)
                elif conversion_type == 'pdf_to_text':
                    text_filename = f'{os.path.splitext(filename)[0]}.txt'
                    convert_pdf_to_text(file_path, os.path.join(app.config['CONVERTED_FOLDER'], text_filename))
                    text_url = url_for('converted_files', filename=text_filename)
                    return render_template('index.html', text_url=text_url)
                elif conversion_type == 'pdf_to_pptx':
                    pptx_filename = f'{os.path.splitext(filename)[0]}.pptx'
                    convert_pdf_to_pptx(file_path, os.path.join(app.config['CONVERTED_FOLDER'], pptx_filename))
                    pptx_url = url_for('converted_files', filename=pptx_filename)
                    return render_template('index.html', pptx_url=pptx_url)
                elif conversion_type == 'pdf_to_docx':
                    docx_filename = f'{os.path.splitext(filename)[0]}.docx'
                    convert_pdf_to_docx(file_path, os.path.join(app.config['CONVERTED_FOLDER'], docx_filename))
                    docx_url = url_for('converted_files', filename=docx_filename)
                    return render_template('index.html', docx_url=docx_url)


                else:
                    return "Unsupported conversion type"
            else:
                if file_extension == '.pptx':
                    convert_ppt_to_pdf(file_path, os.path.join(app.config['CONVERTED_FOLDER'], pdf_filename))
                elif file_extension in (
                '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.pbm', '.pgm', '.ppm', '.xbm', '.webp'):
                    convert_image_to_pdf(file_path, os.path.join(app.config['CONVERTED_FOLDER'], pdf_filename))
                elif file_extension in (
                '.doc', '.docx', '.docm', '.dot', '.dotm', '.dotx', '.odt', '.rtf', '.txt', '.wps', '.xml', '.xps'):
                    convert_doc_to_pdf(file_path, os.path.join(app.config['CONVERTED_FOLDER'], pdf_filename))
                elif file_extension in ('.xlsx', '.xlsm', '.xltx'):
                    convert_excel_to_pdf(file_path, os.path.join(app.config['CONVERTED_FOLDER'], pdf_filename))
                elif file_extension == '.csv':
                    convert_csv_to_pdf(file_path, os.path.join(app.config['CONVERTED_FOLDER'], pdf_filename))
                pdf_url = url_for('uploaded_pdf', filename=pdf_filename)
                return render_template('index.html', pdf_url=pdf_url)
        else:
            return jsonify({"error": "Unsupported file format"})
    return jsonify({"error": "An error occurred"})

# Serving the folders for file conversion
@app.route('/uploads/<filename>')
def uploaded_file(filename):
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename)

@app.route('/converted_pdfs/<filename>')
def uploaded_pdf(filename):
    return send_from_directory(app.config['CONVERTED_FOLDER'], filename, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)
import subprocess
import brotli
import pdfplumber
import pandas as pd
import tabula
from pdfminer.high_level import extract_text
from flask import Flask, render_template, request, redirect, url_for, send_from_directory, jsonify, send_file
import os
import io
import traceback
import shutil
import zipfile
from PIL import Image
from reportlab.lib.pagesizes import letter, landscape, A4
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, PageBreak
from openpyxl import load_workbook
from reportlab.pdfgen import canvas
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, PageBreak
import csv
from PyPDF2 import PdfReader, PdfWriter, PdfFileMerger, PageObject
from werkzeug.utils import secure_filename
from pdf2docx import Converter
import PyPDF2
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_path
from PyPDF2 import PdfReader
from PyPDF2 import PdfMerger
import PyPDF2
# Naming current module
app = Flask(__name__)

# Define upload and converted PDF directories
UPLOAD_FOLDER = 'uploads'
CONVERTED_FOLDER = 'converted_pdfs'
TEMP_FOLDER = 'tmp'
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['CONVERTED_FOLDER'] = CONVERTED_FOLDER
app.config['TEMP_FOLDER'] = TEMP_FOLDER

# Render the home page for file conversion
@app.route('/')
def index():
    return render_template('index.html')

# Function to convert images to PDF
def convert_image_to_pdf(image_path, output_path):
    # A4 size
    pdf_canvas = canvas.Canvas(output_path, pagesize=(595, 842))
    img = Image.open(image_path)
    img_width, img_height = img.size
    scale_factor = min(595 / img_width, 842 / img_height)
    x_offset = (595 - img_width * scale_factor) / 2
    y_offset = (842 - img_height * scale_factor) / 2
    # Draw the image on the PDF
    pdf_canvas.drawInlineImage(image_path, x_offset, y_offset, width=img_width * scale_factor, height=img_height * scale_factor)
    # Save the PDF
    pdf_canvas.save()


# Specify the path to the LibreOffice binary
LIBREOFFICE_PATH = '/usr/bin'

# Set the UNO_PATH environment variable
os.environ['UNO_PATH'] = LIBREOFFICE_PATH

# Function to convert document types to PDF
def convert_doc_to_pdf(document_path, output_path):
    try:
        # Using libreoffice to convert the document to PDF
        subprocess_args = [
            'libreoffice',
            '--headless',
            '--convert-to', 'pdf:writer_pdf_Export',
            '--outdir', os.path.dirname(output_path),
            '--writer_pdf_Export_PageSize', 'A4',
            document_path
        ]
        subprocess.run(subprocess_args, check=True)
        return True
    except subprocess.CalledProcessError as e:
        print(f"Conversion failed: {e}")
        return False


# Function for convering PPT to PDF
def convert_ppt_to_pdf(ppt_path, output_path):
    try:
        # Using libre office for converting ppt to PDF
        env = os.environ.copy()
        env['PDFA1B_OUTDIR'] = os.path.dirname(output_path)
        subprocess_args = [os.path.join(LIBREOFFICE_PATH, 'libreoffice'), '--headless', '--convert-to', 'pdf:writer_pdf_Export', ppt_path]
        subprocess.run(subprocess_args, check=True, env=env)
        return True
    except subprocess.CalledProcessError as e:
        print(f"Conversion failed: {e}")
        return False


# Function to convert CSV to PDF
def convert_csv_to_pdf(csv_path, pdf_path):
    # Reading the CSV file
    data = []
    with open(csv_path, 'r') as csvfile:
        csvreader = csv.reader(csvfile)
        for row in csvreader:
            data.append(row)

    # Create a PDF template using reportlab
    doc = SimpleDocTemplate(pdf_path, pagesize=A4)
    elements = []
    font_name = 'Helvetica-Bold'
    font_size = 8
    columns_per_table = 5
    num_columns = len(data[0])
    for start_col in range(0, num_columns, columns_per_table):
        end_col = start_col + columns_per_table
        table_data = [row[start_col:end_col] for row in data]
        # Generate table with data for PDF
        table = Table(table_data)

        col_widths = [1.5] * len(table_data[0])
        style = TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), (0.9, 0.9, 0.9)),
            ('TEXTCOLOR', (0, 0), (-1, 0), (0, 0, 0)),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), font_name),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), (0.85, 0.85, 0.85)),
            ('GRID', (0, 0), (-1, -1), 1, (0, 0, 0)),
            ('FONTSIZE', (0, 0), (-1, -1), font_size),
        ])
        table.setStyle(style)
        elements.append(table)
        if end_col < num_columns:
            elements.append(PageBreak())
    doc.build(elements)

# Function to convert Excel to PDF
def convert_excel_to_pdf(excel_path, output_path):
    # Create a landscape PDF using the reportlab library
    doc = SimpleDocTemplate(output_path, pagesize=A4)
    elements = []
    workbook = load_workbook(excel_path)
    # Font size and style
    font_name = 'Helvetica-Bold'
    font_size = 12
    # Limiting colunm in pdf
    columns_per_table = 5
    # Reading Sheets from Excel
    for sheet in workbook.sheetnames:
        data = []
        for row in workbook[sheet].iter_rows(values_only=True):
            data.append(row)
        num_columns = len(data[0])
        for start_col in range(0, num_columns, columns_per_table):
            end_col = start_col + columns_per_table
            table_data = [row[start_col:end_col] for row in data]
            # Generate table with data for PDF
            table = Table(table_data)
            col_widths = [1.5] * len(table_data[0])
            style = TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), (0.8, 0.8, 0.8)),
                ('TEXTCOLOR', (0, 0), (-1, 0), (1, 1, 1)),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), font_name),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), (0.95, 0.95, 0.95)),
                ('FONTSIZE', (0, 0), (-1, -1), font_size),
            ])
            table.setStyle(style)
            elements.append(table)
            if end_col < num_columns:
                elements.append(PageBreak())
    doc.build(elements)


# Function converting PDF to image
def convert_pdf_to_image(pdf_path, output_dir):
    images = convert_from_path(pdf_path)
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    image_filenames = []
    for i, image in enumerate(images):
        image_filename = f"page_{i + 1}.png"
        image.save(os.path.join(output_dir, image_filename))
        image_filenames.append(image_filename)
    return image_filenames

#Return the images from directory
@app.route('/converted_images/<filename>')
def converted_images(filename):
    return send_from_directory(app.config['CONVERTED_FOLDER'], filename)

#Converting PDF to Excel file
def convert_pdf_to_xlsx(pdf_path, xlsx_path):
    # Open the PDF file
    with pdfplumber.open(pdf_path) as pdf:
        all_tables = []
        for page in pdf.pages:
            # Extract tables from the page
            tables = page.extract_tables()
            all_tables.extend(tables)
        # Using pandans dataframe for reading table
        df = pd.DataFrame([row for table in all_tables for row in table])
        # Save the dataframe to an excel file
        df.to_excel(xlsx_path, index=False)

#Converting PDF to Csv
def convert_pdf_to_csv(pdf_path, csv_path):
    tabula.convert_into(pdf_path, csv_path, output_format="csv")

@app.route('/converted_files/<filename>')
def converted_files(filename):
    return send_from_directory(app.config['CONVERTED_FOLDER'], filename)

#Converting PDF to text
def convert_pdf_to_text(pdf_path, output_path):
    text = extract_text(pdf_path)
    with open(output_path, "w", encoding="utf-8") as text_file:
        text_file.write(text)


# Function to convert PDF to DOCX
def convert_pdf_to_docx(pdf_path, docx_path):
    try:
        cv = Converter(pdf_path)
        cv.convert(docx_path, start=0, end=None)
        cv.close()
        return True
    except Exception as e:
        print(f"Conversion failed: {e}")
        return False

# Function to convert PDF to PPTX
def convert_pdf_to_pptx(pdf_path, pptx_path):
    # Create a new PowerPoint presentation
    prs = Presentation()
    with open(pdf_path, 'rb') as pdf_file:
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        for page_num, page in enumerate(pdf_reader.pages):
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            shapes = slide.shapes
            title_shape = shapes.title
            if len(shapes.placeholders) > 1:
                body_shape = shapes.placeholders[1]
            else:
                body_shape = shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(5))
            title_shape.text = f"Slide {page_num + 1}"
            tf = body_shape.text_frame
            text = page.extract_text()
            text = text.replace('\n', ' ')
            paragraphs = text.split('\r')
            for paragraph in paragraphs:
                p = tf.add_paragraph()
                p.text = paragraph
    prs.save(pptx_path)


# Rendering the merging function
@app.route('/merge', methods=['GET', 'POST'])
def merge_pdfs():
    if request.method == 'POST':
        pdf_files = request.files.getlist('pdf_files')
        if pdf_files:
            merged_pdf_path = os.path.join(app.config['CONVERTED_FOLDER'], 'merged.pdf')
            merge_pdfs([pdf.filename for pdf in pdf_files], merged_pdf_path)
            merged_pdf_url = url_for('uploaded_pdf', filename='merged.pdf')
            return render_template('merge.html', merged_pdf_url=merged_pdf_url)
        else:
            return "No PDF files to merge."

    return render_template('merge.html')

# Function to merge PDF files
def merge_pdfs(pdf_files, output_path):
    pdf_merger = PdfMerger()
    for pdf_file in pdf_files:
        pdf_merger.append(os.path.join(app.config['UPLOAD_FOLDER'], pdf_file))
    pdf_merger.write(output_path)
    pdf_merger.close()


os.makedirs(app.config['CONVERTED_FOLDER'], exist_ok=True)

# Function to create a zip file
def create_zip_file(zip_filename, files):
    with zipfile.ZipFile(zip_filename, 'w') as zip_file:
        for file in files:
            zip_file.write(file, os.path.basename(file))

# Route to download all split PDFs as a zip file
@app.route('/download_all_zip', methods=['GET', 'POST'])
def download_all_zip():
    if request.method == 'POST':
        pdf_file = request.files['pdf_file']
        if pdf_file:
            output_dir = app.config['CONVERTED_FOLDER']
            os.makedirs(output_dir, exist_ok=True)
            # Define the pdf_reader within the function
            pdf_data = pdf_file.read()
            pdf_reader = PdfReader(io.BytesIO(pdf_data))
            split_pdf_paths = []
            for page_num, page in enumerate(pdf_reader.pages):
                pdf_writer = PdfWriter()
                pdf_writer.add_page(page)
                split_page_path = os.path.join(output_dir, f'page_{page_num + 1}.pdf')
                split_pdf_paths.append(split_page_path)

                with open(split_page_path, 'wb') as output_pdf:
                    pdf_writer.write(output_pdf)

            # Create a zip file
            zip_filename = os.path.join(output_dir, 'split_pdf.zip')
            create_zip_file(zip_filename, split_pdf_paths)

            # Return the zip file and split PDF paths for download
            return render_template('split.html', split_pdf_paths=split_pdf_paths, zip_filename=zip_filename)
    return render_template('split.html')

# Route to download a split PDF
@app.route('/download/<path:filename>')
def download_split_pdf(filename):
    return send_file(filename, as_attachment=True)
    
# Route to download the zip file
@app.route('/download_zip/<path:zip_filename>')
def download_zip(zip_filename):
    return send_file(zip_filename, as_attachment=True)



# Rendering the locking function
@app.route('/lock', methods=['GET', 'POST'])
def lock_pdf():
    if request.method == 'POST':
        pdf_file = request.files['pdf_file']
        password = request.form['password']
        confirm_password = request.form['confirm_password']
        if pdf_file and password == confirm_password:
            output_dir = app.config['CONVERTED_FOLDER']
            os.makedirs(output_dir, exist_ok=True)
            locked_pdf_path = os.path.join(output_dir, secure_filename(pdf_file.filename))
            pdf_file.save(locked_pdf_path)
            set_pdf_password(locked_pdf_path, password)

            locked_pdf_url = url_for('uploaded_pdf', filename=os.path.basename(locked_pdf_path))
            return render_template('lock.html', locked_pdf_url=locked_pdf_url)
        else:
            return "Passwords do not match. Please try again."

    return render_template('lock.html')

# Function to set a password for a PDF
def set_pdf_password(pdf_path, password):
    pdf_writer = PdfWriter()
    with open(pdf_path, 'rb') as pdf_file:
        pdf_reader = PdfReader(pdf_file)
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            pdf_writer.add_page(page)
        pdf_writer.encrypt(password)

    with open(pdf_path, 'wb') as output_pdf:
        pdf_writer.write(output_pdf)


# Rendering the unlocking function
@app.route('/unlock', methods=['GET', 'POST'])
def unlock_pdf():
    if request.method == 'POST':
        pdf_file = request.files['pdf_file']
        password = request.form['password']

        temp_folder = app.config['TEMP_FOLDER']
        temp_pdf_path = os.path.join(temp_folder, secure_filename(pdf_file.filename))
        pdf_file.save(temp_pdf_path)
        if is_pdf_locked(temp_pdf_path, password):
            with open(temp_pdf_path, 'rb') as temp_file:
                pdf_reader = PdfReader(temp_file)
                pdf_reader.decrypt(password)
                pdf_writer = PdfWriter()
                for page_num in range(len(pdf_reader.pages)):
                    pdf_writer.add_page(pdf_reader.pages[page_num])
                unlocked_pdf_path = os.path.join(app.config['CONVERTED_FOLDER'], secure_filename(pdf_file.filename))
                with open(unlocked_pdf_path, 'wb') as unlocked_file:
                    pdf_writer.write(unlocked_file)
            os.remove(temp_pdf_path)

            unlocked_pdf_url = url_for('uploaded_pdf', filename=os.path.basename(unlocked_pdf_path))
            return render_template('unlock.html', unlocked_pdf_url=unlocked_pdf_url)
        else:
            os.remove(temp_pdf_path)
            return "Incorrect password. Please try again."

    return render_template('unlock.html')


# Function to check if a PDF is locked
def is_pdf_locked(pdf_path, password):
    with open(pdf_path, 'rb') as pdf_file:
        pdf_reader = PdfReader(pdf_file)
        return pdf_reader.is_encrypted and pdf_reader.decrypt(password)

os.makedirs('converted_pdfs', exist_ok=True)


def compress_file(input_path, output_path, compression_level):
    try:
        with open(input_path, 'rb') as input_file:
            pdf_reader = PdfReader(input_file)
            pdf_writer = PdfWriter()
            for page_num in range(pdf_reader.getNumPages()):
                page = pdf_reader.getPage(page_num)
                content = page.compressContentStreams()
                if content is not None:
                    # Compress the content using brotli library
                    compressed_content = brotli.compress(content, quality=compression_level)
                    page.__setitem__(PageObject.CONTENTS, compressed_content)
                    pdf_writer.addPage(page)
            with open(output_path, 'wb') as output_file:
                pdf_writer.write(output_file)
        print(f"Compression successful: {output_path}")
    except Exception as e:
        print(f"Error during compression: {e}")

# Rending the compress function
@app.route('/compress', methods=['GET', 'POST'])
def compress():
    import os
    if request.method == 'POST':
        compression_level = int(request.form['compression_level'])
        file = request.files['file']
        if file:
            # Save the uploaded file
            input_path = os.path.join('output', file.filename)
            file.save(input_path)
            # Compress the file
            output_path = os.path.join('output', f'compressed_{file.filename}.zip')
            compress_file(input_path, output_path, compression_level)
            return render_template('compress.html', compressed_file=output_path, os=os)
    return render_template('compress.html', compressed_file=None, os=os)

@app.route('/download_compressed_pdf/<filename>')
def download_compressed_pdf(filename):
    compressed_file_path = os.path.join(app.config['CONVERTED_FOLDER'], filename)
    return send_file(compressed_file_path, as_attachment=True)

# Rendering the delete pages function
@app.route('/delete', methods=['GET', 'POST'])
def delete_pages():
    if request.method == 'POST':
        pdf_file = request.files['pdf_file']
        pages_to_delete = request.form['pages_to_delete']
        # Save the uploaded PDF to a temporary location
        temp_pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], pdf_file.filename)
        pdf_file.save(temp_pdf_path)

        output_pdf_path = os.path.join(app.config['CONVERTED_FOLDER'], 'pages_deleted.pdf')
        delete_pdf_pages(temp_pdf_path, output_pdf_path, pages_to_delete)

        deleted_pdf_url = url_for('uploaded_pdf', filename='pages_deleted.pdf')
        return render_template('delete.html', deleted_pdf_url=deleted_pdf_url)

    return render_template('delete.html')

# Function to delete specific pages from a PDF
def delete_pdf_pages(input_path, output_path, pages_to_delete):
    pdf_reader = PdfReader(input_path)
    pdf_writer = PdfWriter()
    for page_num in range(len(pdf_reader.pages)):
        if str(page_num + 1) not in pages_to_delete.split(','):
            pdf_writer.add_page(pdf_reader.pages[page_num])

    with open(output_path, 'wb') as output_pdf:
        pdf_writer.write(output_pdf)


# Rendering the repair function
@app.route('/repair', methods=['GET', 'POST'])
def repair_pdf():
    if request.method == 'POST':
        pdf_file = request.files['pdf_file']
        temp_pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], pdf_file.filename)
        pdf_file.save(temp_pdf_path)

        output_pdf_path = os.path.join(app.config['CONVERTED_FOLDER'], 'repaired.pdf')
        repair_pdf_file(temp_pdf_path, output_pdf_path)

        repaired_pdf_url = url_for('uploaded_pdf', filename='repaired.pdf')
        return render_template('repair.html', repaired_pdf_url=repaired_pdf_url)

    return render_template('repair.html')

# Function to repair a PDF using PyPDF2
def repair_pdf_file(input_path, output_path):
    pdf_reader = PdfReader(input_path)
    pdf_writer = PdfWriter()

    for page_num in range(pdf_reader.getNumPages()):
        page = pdf_reader.getPage(page_num)
        pdf_writer.addPage(page)

    with open(output_path, 'wb') as output_pdf:
        pdf_writer.write(output_pdf)


# Rendering the rotating function 
@app.route('/rotate', methods=['GET', 'POST'])
def rotate_pdf():
    if request.method == 'POST':
        pdf_file = request.files['pdf_file']
        rotation_direction = request.form.get('rotate_direction')
        if not rotation_direction:
            return "Rotation direction not specified."
        if pdf_file:
            temp_folder = app.config['TEMP_FOLDER']
            temp_pdf_path = os.path.join(temp_folder, secure_filename(pdf_file.filename))
            pdf_file.save(temp_pdf_path)
            rotated_pdf_path = os.path.join(app.config['CONVERTED_FOLDER'], 'rotated.pdf')
            rotate_pdf_pages(temp_pdf_path, rotated_pdf_path, rotation_direction)
            os.remove(temp_pdf_path)
            rotated_pdf_url = url_for('uploaded_pdf', filename='rotated.pdf')
            return render_template('rotate.html', rotated_pdf_url=rotated_pdf_url)
    return render_template('rotate.html')

def rotate_pdf_pages(input_path, output_path, rotation_direction):
    with open(input_path, 'rb') as input_file, open(output_path, 'wb') as output_file:
        pdf_reader = PdfReader(input_file)
        pdf_writer = PdfWriter()
        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            rotation_angle = 90 if rotation_direction == 'right' else -90
            page.rotate(rotation_angle)
            pdf_writer.add_page(page)
        pdf_writer.write(output_file)


# Rendering the converting function for converting supported formats
@app.route('/convert', methods=['POST'])
def convert_file():
    if 'file' not in request.files:
        return redirect(request.url)
    file = request.files['file']
    if file.filename == '':
        return redirect(request.url)
    if file:
        filename = file.filename
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(file_path)
        pdf_filename = f'{os.path.splitext(filename)[0]}.pdf'
        supported_extensions = {
            '.pptx', '.ppt', '.pps', '.ppsx', '.odp',
            '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.pbm', '.pgm', '.ppm', '.xbm', '.webp',
            '.doc', '.docx', '.docm', '.dot', '.dotm', '.dotx', '.odt', '.rtf', '.txt', '.wps', '.xml', '.xps',
            '.xlsx', '.xlsm', '.xltx', '.csv', '.pdf', '.odt', '.ods', '.odp', '.odg', '.odf'
        }
        file_extension = os.path.splitext(filename)[-1].lower()
        if file_extension in supported_extensions:
            if file_extension == '.pdf':
                conversion_type = request.form['conversion_type']
                if conversion_type == 'pdf_to_image':
                    image_filenames = convert_pdf_to_image(file_path, app.config['CONVERTED_FOLDER'])
                    image_urls = [url_for('converted_images', filename=image_filename) for image_filename in
                                  image_filenames]
                    return render_template('index.html', image_urls=image_urls)
                elif conversion_type == 'pdf_to_xlsx':
                    xlsx_filename = f'{os.path.splitext(filename)[0]}.xlsx'
                    convert_pdf_to_xlsx(file_path, os.path.join(app.config['CONVERTED_FOLDER'], xlsx_filename))
                    xlsx_url = url_for('converted_files', filename=xlsx_filename)
                    return render_template('index.html', xlsx_url=xlsx_url)
                elif conversion_type == 'pdf_to_csv':
                    csv_filename = f'{os.path.splitext(filename)[0]}.csv'
                    convert_pdf_to_csv(file_path, os.path.join(app.config['CONVERTED_FOLDER'], csv_filename))
                    csv_url = url_for('converted_files', filename=csv_filename)
                    return render_template('index.html', csv_url=csv_url)
                elif conversion_type == 'pdf_to_text':
                    text_filename = f'{os.path.splitext(filename)[0]}.txt'
                    convert_pdf_to_text(file_path, os.path.join(app.config['CONVERTED_FOLDER'], text_filename))
                    text_url = url_for('converted_files', filename=text_filename)
                    return render_template('index.html', text_url=text_url)
                elif conversion_type == 'pdf_to_pptx':
                    pptx_filename = f'{os.path.splitext(filename)[0]}.pptx'
                    convert_pdf_to_pptx(file_path, os.path.join(app.config['CONVERTED_FOLDER'], pptx_filename))
                    pptx_url = url_for('converted_files', filename=pptx_filename)
                    return render_template('index.html', pptx_url=pptx_url)
                elif conversion_type == 'pdf_to_docx':
                    docx_filename = f'{os.path.splitext(filename)[0]}.docx'
                    convert_pdf_to_docx(file_path, os.path.join(app.config['CONVERTED_FOLDER'], docx_filename))
                    docx_url = url_for('converted_files', filename=docx_filename)
                    return render_template('index.html', docx_url=docx_url)
                else:
                    return "Unsupported conversion type"
            else:
                if file_extension == '.pptx':
                    convert_ppt_to_pdf(file_path, os.path.join(app.config['CONVERTED_FOLDER'], pdf_filename))
                elif file_extension in (
                '.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.pbm', '.pgm', '.ppm', '.xbm', '.webp'):
                    convert_image_to_pdf(file_path, os.path.join(app.config['CONVERTED_FOLDER'], pdf_filename))
                elif file_extension in (
                '.doc', '.docx', '.docm', '.dot', '.dotm', '.dotx', '.odt', '.rtf', '.txt', '.wps', '.xml', '.xps'):
                    convert_doc_to_pdf(file_path, os.path.join(app.config['CONVERTED_FOLDER'], pdf_filename))
                elif file_extension in ('.xlsx', '.xlsm', '.xltx'):
                    convert_excel_to_pdf(file_path, os.path.join(app.config['CONVERTED_FOLDER'], pdf_filename))
                elif file_extension == '.csv':
                    convert_csv_to_pdf(file_path, os.path.join(app.config['CONVERTED_FOLDER'], pdf_filename))
                pdf_url = url_for('uploaded_pdf', filename=pdf_filename)
                return render_template('index.html', pdf_url=pdf_url)
        else:
            return jsonify({"error": "Unsupported file format"})
    return jsonify({"error": "An error occurred"})

# Serving the folders for file conversion
@app.route('/uploads/<filename>')
def uploaded_file(filename):
    return send_from_directory(app.config['UPLOAD_FOLDER'], filename)

@app.route('/converted_pdfs/<filename>')
def uploaded_pdf(filename):
    return send_from_directory(app.config['CONVERTED_FOLDER'], filename, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True)
